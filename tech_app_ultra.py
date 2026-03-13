from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


def create_ultra_design_ppt():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    # Theme Palette
    NAVY_BLACK = RGBColor(3, 7, 18)
    ELECTRIC_BLUE = RGBColor(0, 212, 255)
    DEEP_PURPLE = RGBColor(99, 102, 241)
    CYAN_GLOW = RGBColor(34, 211, 238)

    def apply_blueprint_bg(slide):
        # 1. Dark Base
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = NAVY_BLACK
        bg.line.fill.background()

        # 2. Tech Grid (Vertical data lines)
        for i in range(20):
            x = i * Inches(0.7)
            line = slide.shapes.add_connector(MSO_SHAPE.RECTANGLE, x, 0, x, prs.slide_height)
            line.line.color.rgb = RGBColor(15, 23, 42)
            line.line.width = Pt(0.5)

        # 3. Corner Data Nodes
        for x, y in [(0, 0), (prs.slide_width - Inches(1), 0)]:
            node = slide.shapes.add_shape(MSO_SHAPE.FLOWCHART_CONNECTOR, x, y, Inches(1), Inches(1))
            node.fill.solid()
            node.fill.fore_color.rgb = ELECTRIC_BLUE
            node.line.fill.background()

    def draw_technical_diagram(slide, x_start, y_center):
        """Creates a professional 3-tier node diagram."""
        # Mobile Sender
        s_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x_start, y_center - Inches(0.5), Inches(2), Inches(1)
        )
        s_box.fill.solid()
        s_box.fill.fore_color.rgb = RGBColor(30, 41, 59)
        s_box.line.color.rgb = ELECTRIC_BLUE

        # ESP32 Relay Hub (diamond-ish using a parallelogram as fallback)
        hub = slide.shapes.add_shape(
            MSO_SHAPE.PARALLELOGRAM, x_start + Inches(3.5), y_center - Inches(0.75), Inches(1.5), Inches(1.5)
        )
        hub.fill.solid()
        hub.fill.fore_color.rgb = DEEP_PURPLE
        hub.line.color.rgb = CYAN_GLOW

        # Mobile Receiver
        r_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x_start + Inches(6.5), y_center - Inches(0.5), Inches(2), Inches(1)
        )
        r_box.fill.solid()
        r_box.fill.fore_color.rgb = RGBColor(30, 41, 59)
        r_box.line.color.rgb = ELECTRIC_BLUE

        # Connector Lines (Data Flow)
        c1 = slide.shapes.add_connector(
            MSO_SHAPE.RECTANGLE, x_start + Inches(2), y_center, x_start + Inches(3.5), y_center
        )
        c1.line.color.rgb = CYAN_GLOW
        c1.line.width = Pt(2)
        c2 = slide.shapes.add_connector(
            MSO_SHAPE.RECTANGLE, x_start + Inches(5), y_center, x_start + Inches(6.5), y_center
        )
        c2.line.color.rgb = CYAN_GLOW
        c2.line.width = Pt(2)

    # --- GENERATE SLIDE: SYSTEM FLOW ---
    s_arch = prs.slides.add_slide(prs.slide_layouts[6])
    apply_blueprint_bg(s_arch)
    draw_technical_diagram(s_arch, Inches(2), Inches(4))

    prs.save("PeerLink_Ultra_Design.pptx")


if __name__ == "__main__":
    create_ultra_design_ppt()

