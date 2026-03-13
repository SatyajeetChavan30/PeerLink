from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR


def create_advanced_peerlink_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Enhanced Color Palette
    BG_COLOR = RGBColor(5, 10, 20)
    CYAN = RGBColor(0, 212, 255)
    PURPLE = RGBColor(124, 58, 237)
    WHITE = RGBColor(255, 255, 255)
    GRAY = RGBColor(100, 116, 139)
    DARK_CARD = RGBColor(15, 23, 42)
    BORDER_GLOW = RGBColor(30, 64, 175)

    def apply_background_design(slide):
        """Adds a tech-grid, glowing bottom bar, and corner ornaments."""
        # 1. Solid Background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background()

        # 2. Tech Grid Pattern (Faint Lines)
        grid_color = RGBColor(15, 30, 60)
        grid_spacing = Inches(0.75)
        # Vertical lines
        for i in range(1, int(prs.slide_width / grid_spacing) + 1):
            x = i * grid_spacing
            line = slide.shapes.add_connector(MSO_SHAPE.RECTANGLE, x, 0, x, prs.slide_height)
            line.line.color.rgb = grid_color
            line.line.width = Pt(0.5)
        # Horizontal lines
        for i in range(1, int(prs.slide_height / grid_spacing) + 1):
            y = i * grid_spacing
            line = slide.shapes.add_connector(MSO_SHAPE.RECTANGLE, 0, y, prs.slide_width, y)
            line.line.color.rgb = grid_color
            line.line.width = Pt(0.5)

        # 3. Corner Ornaments (Circuit Nodes)
        def draw_node(x, y):
            circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.1), Inches(0.1))
            circ.fill.solid()
            circ.fill.fore_color.rgb = CYAN
            circ.line.fill.background()

        draw_node(Inches(0.5), Inches(0.5))
        draw_node(prs.slide_width - Inches(0.6), Inches(0.5))

        # 4. Bottom Glow Bar
        glow = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.04), prs.slide_width, Inches(0.04)
        )
        glow.fill.solid()
        glow.fill.fore_color.rgb = CYAN
        glow.line.fill.background()

    def add_styled_text(
        slide,
        text,
        left,
        top,
        width,
        height,
        size,
        color=WHITE,
        bold=False,
        align=PP_ALIGN.LEFT,
        font="Inter",
    ):
        tx = slide.shapes.add_textbox(left, top, width, height)
        tf = tx.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font
        return tx

    def add_tech_card(slide, left, top, width, height, icon, title, body):
        # Card Body
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        rect.fill.solid()
        rect.fill.fore_color.rgb = DARK_CARD
        rect.line.color.rgb = RGBColor(30, 41, 59)
        rect.line.width = Pt(1.5)

        # Header Highlight
        add_styled_text(slide, icon, left + Inches(0.1), top + Inches(0.1), width, Inches(0.5), 24)
        add_styled_text(
            slide,
            title,
            left + Inches(0.2),
            top + Inches(0.6),
            width - Inches(0.4),
            Inches(0.4),
            14,
            CYAN,
            True,
        )
        add_styled_text(
            slide,
            body,
            left + Inches(0.2),
            top + Inches(1.0),
            width - Inches(0.4),
            height - Inches(1.1),
            11,
            GRAY,
        )

    def add_code_window(slide, left, top, width, height, code_text):
        # Window Frame
        frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        frame.fill.solid()
        frame.fill.fore_color.rgb = RGBColor(10, 10, 15)
        frame.line.color.rgb = PURPLE

        # Window Controls (Dots)
        for i, color in enumerate(
            [RGBColor(255, 95, 87), RGBColor(255, 189, 46), RGBColor(39, 201, 63)]
        ):
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, left + Inches(0.15 + (i * 0.2)), top + Inches(0.1), Inches(0.1), Inches(0.1)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = color
            dot.line.fill.background()

        add_styled_text(
            slide,
            code_text,
            left + Inches(0.2),
            top + Inches(0.4),
            width - Inches(0.4),
            height - Inches(0.5),
            10,
            RGBColor(200, 200, 200),
            font="Courier New",
        )

    # ══════════════ SLIDE 1: HERO TITLE ══════════════
    s1 = prs.slides.add_slide(blank_layout)
    apply_background_design(s1)

    # Large Decorative Ring in Background
    ring = s1.shapes.add_shape(
        MSO_SHAPE.OVAL, prs.slide_width / 2 - Inches(2), Inches(1), Inches(4), Inches(4)
    )
    ring.fill.background()
    ring.line.color.rgb = RGBColor(20, 40, 80)
    ring.line.width = Pt(2)

    add_styled_text(s1, "PEERLINK", 0, Inches(2.8), prs.slide_width, Inches(1), 64, WHITE, True, PP_ALIGN.CENTER)
    add_styled_text(
        s1,
        "OFFLINE MESH COMMUNICATION SYSTEM",
        0,
        Inches(3.8),
        prs.slide_width,
        Inches(0.4),
        14,
        CYAN,
        True,
        PP_ALIGN.CENTER,
    )

    # Footer-style details
    add_styled_text(
        s1,
        "Final Year Project | 2026",
        Inches(1),
        prs.slide_height - Inches(0.8),
        Inches(4),
        Inches(0.4),
        11,
        GRAY,
    )

    # ══════════════ SLIDE 4: SYSTEM ARCHITECTURE ══════════════
    s4 = prs.slides.add_slide(blank_layout)
    apply_background_design(s4)
    add_styled_text(
        s4, "SYSTEM ARCHITECTURE", Inches(0.8), Inches(0.6), Inches(5), Inches(0.5), 28, WHITE, True
    )

    # Drawing a more "Flow-y" Diagram
    nodes = [
        (Inches(1), Inches(2.5), "MOBILE SENDER", "📱"),
        (Inches(4.5), Inches(2.5), "ESP32 RELAY", "⚡"),
        (Inches(8), Inches(2.5), "MOBILE RECEIVER", "📲"),
    ]

    for i, (x, y, label, icon) in enumerate(nodes):
        add_tech_card(s4, x, y, Inches(2.5), Inches(1.5), icon, label, "Data Packet Node")
        if i < len(nodes) - 1:
            # Add decorative arrow line
            conn = s4.shapes.add_connector(
                MSO_SHAPE.RECTANGLE, x + Inches(2.5), y + Inches(0.75), x + Inches(3.5), y + Inches(0.75)
            )
            conn.line.color.rgb = CYAN
            conn.line.width = Pt(2)

    # ══════════════ SLIDE 5: FIRMWARE CODE ══════════════
    s5 = prs.slides.add_slide(blank_layout)
    apply_background_design(s5)
    add_styled_text(
        s5, "ESP32 FIRMWARE LOGIC", Inches(0.8), Inches(0.6), Inches(6), Inches(0.5), 28, WHITE, True
    )

    code_snippet = """void loop() {
    int packetSize = UDP.parsePacket();
    if (packetSize) {
        String type = json["type"];
        if (type == "chat") {
            routePacket(json["payload"], json["to"]);
            blinkLED(GPIO_2);
        }
    }
    checkClientTimeouts();
}"""
    add_code_window(s5, Inches(1), Inches(2), Inches(6), Inches(4), code_snippet)

    # Feature list next to code
    features = "• Port: 8888\n• Protocol: UDP\n• Format: JSON\n• Range: ~50m"
    add_styled_text(s5, features, Inches(8), Inches(2.5), Inches(4), Inches(2), 18, CYAN, True)

    # ══════════════ SLIDE 9: TECH STACK ══════════════
    s9 = prs.slides.add_slide(blank_layout)
    apply_background_design(s9)
    add_styled_text(
        s9, "TECHNOLOGY STACK", Inches(0.8), Inches(0.6), Inches(6), Inches(0.5), 28, WHITE, True
    )

    # 3-Column Design
    stacks = [
        ("KOTLIN", "Jetpack Compose\nRoom DB\nCoroutines", Inches(1)),
        ("C++", "Arduino SDK\nWiFi.h\nArduinoJson", Inches(5)),
        ("NETWORK", "UDP Protocol\nLocal Mesh\nSpecifier API", Inches(9)),
    ]

    for name, desc, x in stacks:
        rect = s9.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2), Inches(3.3), Inches(4))
        rect.fill.solid()
        rect.fill.fore_color.rgb = DARK_CARD
        rect.line.color.rgb = CYAN
        add_styled_text(s9, name, x, Inches(2.2), Inches(3.3), Inches(0.5), 20, CYAN, True, PP_ALIGN.CENTER)
        add_styled_text(
            s9,
            desc,
            x + Inches(0.2),
            Inches(3),
            Inches(2.9),
            Inches(2),
            14,
            GRAY,
            False,
            PP_ALIGN.CENTER,
        )

    prs.save("PeerLink_Designer_Presentation.pptx")
    print("Designer Presentation Saved!")


if __name__ == "__main__":
    create_advanced_peerlink_ppt()

