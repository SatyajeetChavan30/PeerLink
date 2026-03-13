from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


def create_tech_app_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    BG_COLOR = RGBColor(5, 10, 20)
    CYAN = RGBColor(0, 212, 255)
    PURPLE = RGBColor(124, 58, 237)
    WHITE = RGBColor(255, 255, 255)
    GRAY = RGBColor(100, 116, 139)
    DARK_CARD = RGBColor(15, 23, 42)
    ELECTRIC_BLUE = CYAN
    DEEP_PURPLE = PURPLE
    CYAN_GLOW = RGBColor(34, 211, 238)

    def apply_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background()

        # simple bottom glow
        glow = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.05), prs.slide_width, Inches(0.05)
        )
        glow.fill.solid()
        glow.fill.fore_color.rgb = CYAN
        glow.line.fill.background()

    def add_text(slide, text, left, top, width, height, size, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
        tx = slide.shapes.add_textbox(left, top, width, height)
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Inter"
        return tx

    def add_hero_title_with_team():
        # HERO TITLE WITH TEAM
        s1 = prs.slides.add_slide(prs.slide_layouts[6])
        apply_bg(s1)

        ring = s1.shapes.add_shape(
            MSO_SHAPE.DONUT, prs.slide_width / 2 - Inches(2), Inches(0.8), Inches(4), Inches(4)
        )
        ring.fill.background()
        ring.line.color.rgb = RGBColor(20, 40, 80)
        ring.line.width = Pt(1.5)

        add_text(
            s1,
            "PEERLINK",
            0,
            Inches(2.2),
            prs.slide_width,
            Inches(1),
            64,
            WHITE,
            True,
            PP_ALIGN.CENTER,
        )
        add_text(
            s1,
            "OFFLINE MESH COMMUNICATION SYSTEM",
            0,
            Inches(3.2),
            prs.slide_width,
            Inches(0.4),
            14,
            ELECTRIC_BLUE,
            True,
            PP_ALIGN.CENTER,
        )

        add_text(
            s1,
            "Final Year Project | 2026",
            0,
            Inches(3.6),
            prs.slide_width,
            Inches(0.4),
            12,
            GRAY,
            False,
            PP_ALIGN.CENTER,
        )

        add_text(
            s1,
            "CORE DEVELOPMENT TEAM",
            Inches(1),
            Inches(4.8),
            Inches(11.333),
            Inches(0.4),
            10,
            DEEP_PURPLE,
            True,
            PP_ALIGN.CENTER,
        )

        team_members = [
            ("Alex Mercer", "Lead Android Developer"),
            ("Sarah Chen", "Firmware & ESP32"),
            ("Marcus Johnson", "Network Protocol"),
            ("Priya Patel", "UI/UX & DB Architecture"),
        ]

        card_width = Inches(2.5)
        card_height = Inches(0.8)
        spacing = Inches(0.3)
        start_x = Inches(1.2)
        y_pos = Inches(5.2)

        for i, (name, role) in enumerate(team_members):
            x_pos = start_x + (i * (card_width + spacing))

            badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, y_pos, card_width, card_height)
            badge.fill.solid()
            badge.fill.fore_color.rgb = RGBColor(10, 15, 30)
            badge.line.color.rgb = RGBColor(30, 60, 100)
            badge.line.width = Pt(1)

            dot = s1.shapes.add_shape(
                MSO_SHAPE.OVAL, x_pos + Inches(0.1), y_pos + Inches(0.35), Inches(0.1), Inches(0.1)
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = ELECTRIC_BLUE if i % 2 == 0 else CYAN_GLOW
            dot.line.fill.background()

            add_text(
                s1,
                name,
                x_pos + Inches(0.3),
                y_pos + Inches(0.15),
                card_width - Inches(0.3),
                Inches(0.3),
                12,
                WHITE,
                True,
                PP_ALIGN.LEFT,
            )
            add_text(
                s1,
                role,
                x_pos + Inches(0.3),
                y_pos + Inches(0.4),
                card_width - Inches(0.3),
                Inches(0.3),
                9,
                GRAY,
                False,
                PP_ALIGN.LEFT,
            )

    def add_technical_approach_slide():
        s_tech = prs.slides.add_slide(prs.slide_layouts[6])
        apply_bg(s_tech)

        add_text(
            s_tech,
            "TECHNICAL APPROACH",
            Inches(0.8),
            Inches(0.6),
            Inches(6),
            Inches(0.6),
            32,
            bold=True,
        )

        steps = [
            (
                "01. Data Encapsulation",
                "Messages are wrapped in JSON objects containing 'type', 'payload', and 'to' fields for structured routing.",
                Inches(1.5),
            ),
            (
                "02. Transport Layer",
                "Utilizes UDP Port 8888 for low-latency, connectionless broadcasting across the mesh relay.",
                Inches(3.2),
            ),
            (
                "03. Asynchronous Handling",
                "Mobile app employs Kotlin Coroutines to manage concurrent network I/O without blocking the UI.",
                Inches(4.9),
            ),
        ]

        for title, desc, y_pos in steps:
            box = s_tech.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), y_pos, Inches(11), Inches(1.4))
            box.fill.solid()
            box.fill.fore_color.rgb = DARK_CARD
            box.line.color.rgb = PURPLE
            box.line.width = Pt(1)

            add_text(s_tech, title, Inches(1.3), y_pos + Inches(0.2), Inches(5), Inches(0.4), 18, CYAN, True)
            add_text(s_tech, desc, Inches(1.3), y_pos + Inches(0.6), Inches(10), Inches(0.6), 13, GRAY)

    add_hero_title_with_team()
    add_technical_approach_slide()
    prs.save("Tech_App.pptx")


if __name__ == "__main__":
    create_tech_app_ppt()

