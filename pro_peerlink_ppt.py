from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
import random


def create_pro_peerlink_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Theme Colors
    BG_COLOR = RGBColor(5, 10, 20)
    CYAN = RGBColor(0, 212, 255)
    PURPLE = RGBColor(124, 58, 237)
    WHITE = RGBColor(255, 255, 255)
    GRAY = RGBColor(100, 116, 139)
    DARK_CARD = RGBColor(15, 23, 42)

    def apply_advanced_background(slide):
        """Adds binary strings, pulse waves, and hexagonal mesh to the background."""
        # 1. Base Background
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.background()

        # 2. Hexagonal Mesh Pattern (Ghosted)
        hex_size = Inches(0.4)
        for _ in range(5):  # Add a few scattered hex clusters
            hx = slide.shapes.add_shape(
                MSO_SHAPE.HEXAGON, int(random.uniform(0, 12) * 914400), int(random.uniform(0, 6) * 914400), hex_size, hex_size
            )
            hx.fill.solid()
            hx.fill.fore_color.rgb = RGBColor(10, 25, 50)
            hx.line.color.rgb = RGBColor(20, 50, 80)
            hx.rotation = 30

        # 3. Binary Code Overlays (Low Opacity effect)
        binary_str = "01101001 01101110 01101011"  # "ink" in binary
        for _ in range(3):
            tx = slide.shapes.add_textbox(
                int(random.uniform(0, 10) * 914400), int(random.uniform(0, 7) * 914400), Inches(3), Inches(0.5)
            )
            p = tx.text_frame.paragraphs[0]
            p.text = binary_str
            p.font.size = Pt(10)
            p.font.name = "Courier New"
            p.font.color.rgb = RGBColor(15, 35, 65)

        # 4. Signal Pulse (Concentric Arcs)
        for r in [1, 1.2, 1.4]:
            arc = slide.shapes.add_shape(MSO_SHAPE.ARC, Inches(-0.5), Inches(-0.5), Inches(r), Inches(r))
            arc.line.color.rgb = RGBColor(20, 60, 100)
            arc.line.width = Pt(1)

        # 5. Bottom Glow Accent
        glow = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.05), prs.slide_width, Inches(0.05)
        )
        glow.fill.solid()
        glow.fill.fore_color.rgb = CYAN
        glow.line.fill.background()

    def add_styled_text(
        slide,
        text,
        left,
        top,
        width,
        height,
        size,
        color=WHITE,
        bold=False,
        align=PP_ALIGN.LEFT,
    ):
        tx = slide.shapes.add_textbox(left, top, width, height)
        tf = tx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Inter"
        return tx

    # --- SLIDE 1: HERO TITLE ---
    s1 = prs.slides.add_slide(blank_layout)
    apply_advanced_background(s1)
    add_styled_text(s1, "PEERLINK", 0, Inches(3), prs.slide_width, Inches(1), 60, WHITE, True, PP_ALIGN.CENTER)
    add_styled_text(
        s1,
        "OFFLINE MESH COMMUNICATION SYSTEM",
        0,
        Inches(4),
        prs.slide_width,
        Inches(0.5),
        16,
        CYAN,
        True,
        PP_ALIGN.CENTER,
    )
    add_styled_text(
        s1,
        "Final Year Project | 2026",
        Inches(1),
        prs.slide_height - Inches(1),
        Inches(4),
        Inches(0.5),
        12,
        GRAY,
    )

    # --- NEW SLIDE: LITERATURE SURVEY ---
    s_ls = prs.slides.add_slide(blank_layout)
    apply_advanced_background(s_ls)
    add_styled_text(
        s_ls, "LITERATURE SURVEY", Inches(0.8), Inches(0.6), Inches(6), Inches(0.6), 32, WHITE, True
    )

    surveys = [
        (
            "FireChat / Bridgefy",
            "Uses Bluetooth/Wi-Fi Direct. Proprioritary and often requires initial internet to sync keys.",
        ),
        ("LoRa Mesh Networks", "Great range but extremely low bandwidth. Cannot support real-time chat UI effectively."),
        ("Zigbee Systems", "Low power but requires specialized hardware not present in modern smartphones."),
        ("PeerLink (Ours)", "Uses ESP32 as an ubiquitous Wi-Fi bridge. High bandwidth, zero-config, and open source."),
    ]

    for i, (title, desc) in enumerate(surveys):
        y_pos = Inches(1.8 + (i * 1.3))
        box = s_ls.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), y_pos, Inches(11), Inches(1.1))
        box.fill.solid()
        box.fill.fore_color.rgb = DARK_CARD
        box.line.color.rgb = PURPLE
        add_styled_text(s_ls, title, Inches(1.2), y_pos + Inches(0.1), Inches(4), Inches(0.4), 16, CYAN, True)
        add_styled_text(s_ls, desc, Inches(1.2), y_pos + Inches(0.5), Inches(10), Inches(0.5), 12, GRAY)

    # --- SLIDE: SYSTEM ARCHITECTURE ---
    s4 = prs.slides.add_slide(blank_layout)
    apply_advanced_background(s4)
    add_styled_text(
        s4, "SYSTEM ARCHITECTURE", Inches(0.8), Inches(0.6), Inches(6), Inches(0.6), 32, WHITE, True
    )
    arch_nodes = [
        (Inches(1.5), "MOBILE SENDER", "📱"),
        (Inches(5.5), "ESP32 RELAY", "⚡"),
        (Inches(9.5), "MOBILE RECEIVER", "📲"),
    ]
    for x, label, icon in arch_nodes:
        add_styled_text(s4, icon, x, Inches(2.5), Inches(2), Inches(1), 40, align=PP_ALIGN.CENTER)
        add_styled_text(s4, label, x, Inches(3.5), Inches(2), Inches(0.5), 14, CYAN, True, PP_ALIGN.CENTER)

    # --- SLIDE: TECHNOLOGY STACK ---
    s9 = prs.slides.add_slide(blank_layout)
    apply_advanced_background(s9)
    add_styled_text(
        s9, "TECHNOLOGY STACK", Inches(0.8), Inches(0.6), Inches(6), Inches(0.6), 32, WHITE, True
    )

    stacks = [
        ("KOTLIN", "Jetpack Compose\nRoom DB\nCoroutines", Inches(1)),
        ("C++", "Arduino SDK\nWiFi.h\nArduinoJson", Inches(5)),
        ("NETWORK", "UDP Protocol\nLocal Mesh\nSpecifier API", Inches(9)),
    ]
    for name, tech, x in stacks:
        rect = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2), Inches(3.3), Inches(4))
        rect.fill.solid()
        rect.fill.fore_color.rgb = DARK_CARD
        rect.line.color.rgb = CYAN
        add_styled_text(s9, name, x, Inches(2.2), Inches(3.3), Inches(0.5), 20, CYAN, True, PP_ALIGN.CENTER)
        add_styled_text(s9, tech, x, Inches(3), Inches(3.3), Inches(2.5), 14, WHITE, False, PP_ALIGN.CENTER)

    # add multi-hop routing slide
    add_multi_hop_mesh_slide(prs, apply_advanced_background, add_styled_text)

    prs.save("PeerLink_Full_Presentation.pptx")


def add_multi_hop_mesh_slide(prs, apply_bg, add_text):
    CYAN = RGBColor(0, 212, 255)
    PURPLE = RGBColor(124, 58, 237)
    GRAY = RGBColor(100, 116, 139)
    WHITE = RGBColor(255, 255, 255)
    DARK_CARD = RGBColor(15, 23, 42)
    HIGHLIGHT = RGBColor(57, 255, 20)

    s_mesh = prs.slides.add_slide(prs.slide_layouts[6])
    apply_bg(s_mesh)

    add_text(
        s_mesh,
        "MULTI-HOP ROUTING",
        Inches(0.8),
        Inches(0.6),
        Inches(6),
        Inches(0.6),
        32,
        WHITE,
        True,
    )
    add_text(
        s_mesh,
        "ESP-TO-ESP PACKET FORWARDING",
        Inches(0.8),
        Inches(1.1),
        Inches(6),
        Inches(0.4),
        14,
        CYAN,
        True,
    )

    nodes = {
        "Sender": (Inches(1.0), Inches(3.5), "📱 Sender", False, True),
        "ESP1": (Inches(4.0), Inches(2.0), "⚡ ESP32\nNode A", True, True),
        "ESP2": (Inches(4.0), Inches(5.0), "⚡ ESP32\nNode B", True, False),
        "ESP3": (Inches(7.0), Inches(3.5), "⚡ ESP32\nNode C", True, True),
        "Receiver": (Inches(10.0), Inches(3.5), "📲 Receiver", False, True),
    }

    def draw_line(start_key, end_key, is_highlighted=False):
        start_x, start_y = nodes[start_key][0], nodes[start_key][1]
        end_x, end_y = nodes[end_key][0], nodes[end_key][1]

        offset_start = Inches(0.5) if nodes[start_key][3] else Inches(0.75)
        offset_end = Inches(0.5) if nodes[end_key][3] else Inches(0.4)

        conn = s_mesh.shapes.add_connector(
            MSO_SHAPE.RECTANGLE,
            start_x + offset_start,
            start_y + offset_start,
            end_x + offset_end,
            end_y + offset_end,
        )
        if is_highlighted:
            conn.line.color.rgb = HIGHLIGHT
            conn.line.width = Pt(3.5)
            add_text(
                s_mesh,
                "DATA",
                start_x + (end_x - start_x) / 2,
                start_y + (end_y - start_y) / 2 - Inches(0.2),
                Inches(1),
                Inches(0.3),
                10,
                HIGHLIGHT,
                True,
                PP_ALIGN.CENTER,
            )
        else:
            conn.line.color.rgb = PURPLE
            conn.line.width = Pt(1.5)
            conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH

    draw_line("Sender", "ESP1", True)
    draw_line("Sender", "ESP2", False)
    draw_line("ESP1", "ESP3", True)
    draw_line("ESP2", "ESP3", False)
    draw_line("ESP3", "Receiver", True)

    for key, (x, y, label, is_esp, is_active) in nodes.items():
        shape_type = MSO_SHAPE.HEXAGON if is_esp else MSO_SHAPE.ROUNDED_RECTANGLE
        width = Inches(1.2) if is_esp else Inches(1.5)
        height = Inches(1.2) if is_esp else Inches(0.8)

        node_shape = s_mesh.shapes.add_shape(shape_type, x, y, width, height)
        node_shape.fill.solid()
        node_shape.fill.fore_color.rgb = DARK_CARD

        node_shape.line.color.rgb = HIGHLIGHT if is_active else CYAN
        node_shape.line.width = Pt(2.5) if is_active else Pt(1.5)

        text_y = y + Inches(0.3) if is_esp else y + Inches(0.15)
        add_text(s_mesh, label, x, text_y, width, height, 11, WHITE, True, PP_ALIGN.CENTER)

    exp_box = s_mesh.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(6.0), Inches(11.4), Inches(1.0)
    )
    exp_box.fill.solid()
    exp_box.fill.fore_color.rgb = DARK_CARD
    exp_box.line.color.rgb = PURPLE

    explanation = (
        "When the Receiver is out of direct range, the protocol dynamically routes the payload through intermediate "
        "ESP32 nodes (Node A → Node C). Inactive nodes (Node B) remain on standby to self-heal the network if a "
        "primary node fails."
    )
    add_text(s_mesh, explanation, Inches(1.0), Inches(6.15), Inches(11.0), Inches(0.8), 12, GRAY)


if __name__ == "__main__":
    create_pro_peerlink_ppt()
