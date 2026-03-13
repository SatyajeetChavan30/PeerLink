from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


def create_peerlink_presentation():
    prs = Presentation()
    # Set to Widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Theme Colors
    BG_COLOR = RGBColor(5, 10, 20)
    CYAN = RGBColor(0, 212, 255)
    PURPLE = RGBColor(124, 58, 237)
    WHITE = RGBColor(255, 255, 255)
    GRAY = RGBColor(180, 180, 180)
    CARD_BG = RGBColor(20, 30, 50)
    CARD_BORDER = RGBColor(0, 100, 130)

    # --- Helper Functions ---
    def set_background(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = BG_COLOR
        bg.line.fill.solid()
        bg.line.fill.fore_color.rgb = BG_COLOR
        return bg

    def add_text(
        slide,
        text,
        left,
        top,
        width,
        height,
        font_size,
        color=WHITE,
        bold=False,
        align=PP_ALIGN.LEFT,
        font_name="Arial",
    ):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = align
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        return txBox

    def add_header(slide, slide_num, tag, title):
        set_background(slide)
        # Slide Number
        add_text(
            slide,
            f"{slide_num:02d} / 11",
            Inches(11.5),
            Inches(0.3),
            Inches(1.5),
            Inches(0.5),
            12,
            CYAN,
            align=PP_ALIGN.RIGHT,
            font_name="Courier New",
        )
        # Tag
        add_text(
            slide,
            tag.upper(),
            Inches(1),
            Inches(0.6),
            Inches(11.3),
            Inches(0.4),
            11,
            CYAN,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        # Title
        add_text(
            slide,
            title,
            Inches(1),
            Inches(1),
            Inches(11.3),
            Inches(0.8),
            36,
            WHITE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    def add_card(slide, left, top, width, height, icon, title, text):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = CARD_BORDER

        # Icon
        add_text(slide, icon, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.6), 28)
        # Title
        add_text(
            slide,
            title,
            left + Inches(0.2),
            top + Inches(0.7),
            width - Inches(0.4),
            Inches(0.4),
            14,
            CYAN,
            bold=True,
        )
        # Text
        add_text(
            slide,
            text,
            left + Inches(0.2),
            top + Inches(1.1),
            width - Inches(0.4),
            height - Inches(1.1),
            12,
            GRAY,
        )

    def add_pill(slide, left, top, width, text, color):
        pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, Inches(0.4))
        pill.fill.solid()
        pill.fill.fore_color.rgb = BG_COLOR
        pill.line.color.rgb = color
        add_text(
            slide,
            text,
            left,
            top + Inches(0.05),
            width,
            Inches(0.3),
            11,
            color,
            bold=True,
            align=PP_ALIGN.CENTER,
        )

    def add_flow_node(slide, left, top, width, height, icon, title, sub, color=CYAN):
        node = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        node.fill.solid()
        node.fill.fore_color.rgb = RGBColor(10, 20, 35)
        node.line.color.rgb = color
        add_text(slide, icon, left, top + Inches(0.1), width, Inches(0.5), 24, align=PP_ALIGN.CENTER)
        add_text(
            slide,
            title,
            left,
            top + Inches(0.6),
            width,
            Inches(0.3),
            12,
            color,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            sub,
            left,
            top + Inches(0.9),
            width,
            Inches(0.3),
            10,
            GRAY,
            align=PP_ALIGN.CENTER,
        )

    def add_code_block(slide, left, top, width, height, code):
        block = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        block.fill.solid()
        block.fill.fore_color.rgb = RGBColor(0, 0, 0)
        block.line.color.rgb = CARD_BORDER
        add_text(
            slide,
            code,
            left + Inches(0.2),
            top + Inches(0.2),
            width - Inches(0.4),
            height - Inches(0.4),
            11,
            GRAY,
            font_name="Courier New",
        )

    def add_stat(slide, left, top, width, num, label):
        add_text(
            slide,
            num,
            left,
            top,
            width,
            Inches(0.8),
            44,
            CYAN,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            label,
            left,
            top + Inches(0.8),
            width,
            Inches(0.4),
            12,
            GRAY,
            align=PP_ALIGN.CENTER,
        )

    # ══════════════ SLIDE 1: Title ══════════════
    s1 = prs.slides.add_slide(blank_layout)
    set_background(s1)
    add_text(
        s1,
        "01 / 11",
        Inches(11.5),
        Inches(0.3),
        Inches(1.5),
        Inches(0.5),
        12,
        CYAN,
        align=PP_ALIGN.RIGHT,
        font_name="Courier New",
    )
    add_text(s1, "📡", Inches(6.1), Inches(1.5), Inches(1.2), Inches(1.2), 48, align=PP_ALIGN.CENTER)
    add_text(
        s1,
        "FINAL YEAR PROJECT PRESENTATION",
        Inches(1),
        Inches(2.8),
        Inches(11.3),
        Inches(0.5),
        11,
        CYAN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s1,
        "PeerLink",
        Inches(1),
        Inches(3.2),
        Inches(11.3),
        Inches(1),
        54,
        WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s1,
        "An offline, internet-free peer-to-peer messaging system using an ESP32 microcontroller "
        "as a portable Wi-Fi relay hub for Android devices.",
        Inches(2.6),
        Inches(4.5),
        Inches(8),
        Inches(1),
        14,
        GRAY,
        align=PP_ALIGN.CENTER,
    )

    add_pill(s1, Inches(3.5), Inches(5.8), Inches(1.5), "Android", CYAN)
    add_pill(s1, Inches(5.2), Inches(5.8), Inches(1.5), "ESP32", RGBColor(251, 146, 60))
    add_pill(s1, Inches(6.9), Inches(5.8), Inches(1.5), "Jetpack", PURPLE)
    add_pill(s1, Inches(8.6), Inches(5.8), Inches(1.5), "UDP Mesh", RGBColor(52, 211, 153))

    # ══════════════ SLIDE 2: Problem Statement ══════════════
    s2 = prs.slides.add_slide(blank_layout)
    add_header(s2, 2, "Motivation", "The Problem")
    add_card(
        s2,
        Inches(1.5),
        Inches(2.5),
        Inches(3.2),
        Inches(2.5),
        "🌐",
        "Internet Dependency",
        "All popular messaging apps rely on cloud servers and active internet connectivity to exchange messages.",
    )
    add_card(
        s2,
        Inches(5.0),
        Inches(2.5),
        Inches(3.2),
        Inches(2.5),
        "🚫",
        "Dead Zones",
        "Remote areas, disaster zones, underground venues, and events often have no cellular signal or Wi-Fi infrastructure.",
    )
    add_card(
        s2,
        Inches(8.5),
        Inches(2.5),
        Inches(3.2),
        Inches(2.5),
        "🔒",
        "Privacy & Surveillance",
        "Server-based messaging exposes messages to third-party cloud providers, creating privacy and security concerns.",
    )

    hl = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.5), Inches(5.5), Inches(8.3), Inches(0.8))
    hl.fill.solid()
    hl.fill.fore_color.rgb = RGBColor(10, 30, 45)
    hl.line.color.rgb = CYAN
    add_text(
        s2,
        "💡 The Need: A simple, infrastructure-free communication method that works with just a small device and smartphones.",
        Inches(2.7),
        Inches(5.7),
        Inches(7.9),
        Inches(0.6),
        13,
        WHITE,
    )

    # ══════════════ SLIDE 3: Solution ══════════════
    s3 = prs.slides.add_slide(blank_layout)
    add_header(s3, 3, "Solution", "Introducing PeerLink")

    left_text = (
        "PeerLink transforms an ESP32 microcontroller into a portable Wi-Fi relay hub. Nearby Android phones connect "
        "to it and can chat with each other — completely offline, no SIM card, no internet, no server.\n\n"
        "• Works anywhere — even in remote areas\n"
        "• No registration, no accounts, no cloud\n"
        "• ESP32 costs under ₹400 and runs on a phone charger\n"
        "• End-to-end local network"
    )
    add_text(s3, left_text, Inches(1), Inches(2.5), Inches(6), Inches(4), 14, GRAY)

    add_card(
        s3,
        Inches(7.5),
        Inches(2.5),
        Inches(4.5),
        Inches(3.5),
        "📡",
        "ESP32 Acts As:",
        "\n• Wi-Fi Access Point (hotspot)\n• UDP packet relay / router\n• Client registry & tracker\n• Presence broadcaster",
    )

    # ══════════════ SLIDE 4: Architecture ══════════════
    s4 = prs.slides.add_slide(blank_layout)
    add_header(s4, 4, "System Design", "Architecture Overview")

    add_flow_node(s4, Inches(1.5), Inches(2.2), Inches(2), Inches(1.5), "📱", "Android App", "Sender", CYAN)
    add_text(s4, "→", Inches(3.6), Inches(2.7), Inches(0.4), Inches(0.5), 24, CYAN)
    add_flow_node(s4, Inches(4.1), Inches(2.2), Inches(2.2), Inches(1.5), "⚡", "WifiNetworkSpecifier", "Auto-connects", CYAN)
    add_text(s4, "→", Inches(6.4), Inches(2.7), Inches(0.4), Inches(0.5), 24, CYAN)
    add_flow_node(s4, Inches(6.9), Inches(2.2), Inches(2), Inches(1.5), "🔧", "ESP32 Relay", "Routes port 8888", PURPLE)
    add_text(s4, "→", Inches(9.0), Inches(2.7), Inches(0.4), Inches(0.5), 24, CYAN)
    add_flow_node(s4, Inches(9.5), Inches(2.2), Inches(2), Inches(1.5), "📲", "Android App", "Receiver", CYAN)

    add_card(
        s4,
        Inches(1),
        Inches(4.5),
        Inches(2.6),
        Inches(1.5),
        "",
        "EspConnector",
        "Auto-joins PeerLink Wi-Fi using APIs",
    )
    add_card(
        s4,
        Inches(3.8),
        Inches(4.5),
        Inches(2.6),
        Inches(1.5),
        "",
        "UdpChat",
        "Sends/receives JSON packets on UDP 8888",
    )
    add_card(
        s4,
        Inches(6.6),
        Inches(4.5),
        Inches(2.6),
        Inches(1.5),
        "",
        "MeshService",
        "Android Foreground Service orchestration",
    )
    add_card(
        s4,
        Inches(9.4),
        Inches(4.5),
        Inches(2.6),
        Inches(1.5),
        "",
        "Room DB",
        "Persists chat history locally",
    )

    # ══════════════ SLIDE 5: Firmware ══════════════
    s5 = prs.slides.add_slide(blank_layout)
    add_header(s5, 5, "Hardware Layer", "ESP32 Relay Firmware")

    bullets = (
        "• Wi-Fi AP Mode — Broadcasts SSID PeerLink_XXXX\n"
        "• UDP Server on port 8888 — Listens and routes\n"
        "• Client Registry — Tracks all connected devices\n"
        "• Presence Broadcast — Announces relay every 2s\n"
        "• Client Timeout — Removes inactive devices (30s)\n"
        "• LED Indicator — GPIO 2 blinks on activity"
    )
    add_text(s5, bullets, Inches(1), Inches(2.5), Inches(5.5), Inches(4), 14, GRAY)

    code = (
        '// Message routing logic (simplified)\n'
        'if (type == "chat") {\n'
        '  if (to == "BROADCAST") {\n'
        "    for (auto &[id, info] : clients)\n"
        "      forwardPacket(json, info.ip);\n"
        "  } else {\n"
        "    forwardPacket(json, clients[to].ip);\n"
        "  }\n"
        '} else if (type == "get_clients") {\n'
        "  sendClientList(senderIP);\n"
        "}"
    )
    add_code_block(s5, Inches(6.8), Inches(2.5), Inches(5.5), Inches(3.0), code)

    add_pill(s5, Inches(6.8), Inches(5.8), Inches(1.5), "Arduino C++", RGBColor(251, 146, 60))
    add_pill(s5, Inches(8.5), Inches(5.8), Inches(1.5), "ArduinoJson", CYAN)

    # ══════════════ SLIDE 6: Android App ══════════════
    s6 = prs.slides.add_slide(blank_layout)
    add_header(s6, 6, "Android Layer", "Android Application")
    add_card(
        s6,
        Inches(1.5),
        Inches(2.5),
        Inches(4.8),
        Inches(2.0),
        "🔌",
        "EspConnector",
        "Uses Android's WifiNetworkSpecifier API to connect to hotspot networks via Kotlin callbackFlow.",
    )
    add_card(
        s6,
        Inches(6.8),
        Inches(2.5),
        Inches(4.8),
        Inches(2.0),
        "📦",
        "UdpChat",
        "Manages DatagramSocket on port 8888 for UdpPacket objects. Uses Coroutines & MutableSharedFlow.",
    )
    add_card(
        s6,
        Inches(1.5),
        Inches(4.8),
        Inches(4.8),
        Inches(2.0),
        "⚙️",
        "MeshService",
        "Foreground Service staying alive in background. Orchestrates DB and state flows.",
    )
    add_card(
        s6,
        Inches(6.8),
        Inches(4.8),
        Inches(4.8),
        Inches(2.0),
        "🗄️",
        "Room Database",
        "Persists MessageEntity records locally. UI observes via a Flow<List<MessageEntity>>.",
    )

    # ══════════════ SLIDE 7: App Screens ══════════════
    s7 = prs.slides.add_slide(blank_layout)
    add_header(s7, 7, "User Interface", "App Screen Flow")

    add_flow_node(s7, Inches(0.5), Inches(2.2), Inches(2.2), Inches(1.5), "✨", "Splash Screen", "Brand intro")
    add_text(s7, "→", Inches(2.8), Inches(2.7), Inches(0.3), Inches(0.5), 20, CYAN)
    add_flow_node(s7, Inches(3.1), Inches(2.2), Inches(2.2), Inches(1.5), "🔑", "Permission", "Wi-Fi, Location")
    add_text(s7, "→", Inches(5.4), Inches(2.7), Inches(0.3), Inches(0.5), 20, CYAN)
    add_flow_node(s7, Inches(5.7), Inches(2.2), Inches(2.2), Inches(1.5), "📡", "Discovery", "Connecting")
    add_text(s7, "→", Inches(8.0), Inches(2.7), Inches(0.3), Inches(0.5), 20, CYAN)
    add_flow_node(s7, Inches(8.3), Inches(2.2), Inches(2.2), Inches(1.5), "👥", "Device List", "Shows peers")
    add_text(s7, "→", Inches(10.6), Inches(2.7), Inches(0.3), Inches(0.5), 20, CYAN)
    add_flow_node(s7, Inches(10.9), Inches(2.2), Inches(2.2), Inches(1.5), "💬", "Chat Screen", "Messaging")

    add_card(
        s7,
        Inches(1.5),
        Inches(4.5),
        Inches(3.2),
        Inches(1.8),
        "",
        "Jetpack Compose",
        "Declarative UI with Material 3, reactive state.",
    )
    add_card(
        s7,
        Inches(5.0),
        Inches(4.5),
        Inches(3.2),
        Inches(1.8),
        "",
        "Navigation",
        "Type-safe navigation with NavHost.",
    )
    add_card(
        s7,
        Inches(8.5),
        Inches(4.5),
        Inches(3.2),
        Inches(1.8),
        "",
        "Dark Theme",
        "Black background with white bubbles.",
    )

    # ══════════════ SLIDE 8: Data Flow ══════════════
    s8 = prs.slides.add_slide(blank_layout)
    add_header(s8, 8, "Deep Dive", "Message Data Flow")

    code2 = (
        '// 1. User taps Send → ViewModel → MeshService\n'
        'fun sendMessage(toPeerId: String, content: String) {\n'
        '    val packet = UdpPacket(type="chat", from=myDeviceId, to=toPeerId, payload=content)\n'
        "    database.messageDao().insertMessage(localMsg)  // Save locally first\n"
        "    udpChat.sendMessage(packet)               // Send over UDP → ESP32 router\n"
        "}\n\n"
        "// 2. ESP32 receives packet → looks up target device IP → forwards\n"
        "// 3. Receiver's UdpChat emits via SharedFlow → MeshService handles\n"
        "when (packet.type) {\n"
        '    "chat" -> database.messageDao().insertMessage(incomingMsg)\n'
        '    "clients_list" -> _serviceState.update { it.copy(peers = parsedPeers) }\n'
        "}\n"
        "// 4. Room DB change → Flow emits → Compose UI auto-recomposes"
    )
    add_code_block(s8, Inches(1.5), Inches(2.5), Inches(10.3), Inches(3.5), code2)

    add_pill(s8, Inches(3.5), Inches(6.3), Inches(1.5), "Coroutines", CYAN)
    add_pill(s8, Inches(5.2), Inches(6.3), Inches(1.5), "StateFlow", PURPLE)
    add_pill(s8, Inches(6.9), Inches(6.3), Inches(1.5), "Room Flow", RGBColor(52, 211, 153))

    # ══════════════ SLIDE 9: Tech Stack ══════════════
    s9 = prs.slides.add_slide(blank_layout)
    add_header(s9, 9, "Technology", "Technology Stack")

    add_card(
        s9,
        Inches(1.5),
        Inches(2.5),
        Inches(3.2),
        Inches(2.5),
        "📱",
        "Android / Kotlin",
        "Jetpack Compose · Navigation · ViewModel · Room DB · Foreground Services",
    )
    add_card(
        s9,
        Inches(5.0),
        Inches(2.5),
        Inches(3.2),
        Inches(2.5),
        "⚡",
        "Embedded / C++",
        "ESP32 Arduino SDK · WiFi.h · WiFiUdp.h · ArduinoJson",
    )
    add_card(
        s9,
        Inches(8.5),
        Inches(2.5),
        Inches(3.2),
        Inches(2.5),
        "🌐",
        "Networking",
        "WifiNetworkSpecifier API · UDP Datagrams · Port 8888 · Broadcast routing",
    )

    add_stat(s9, Inches(1.5), Inches(5.5), Inches(2.2), "0", "Internet Required")
    add_stat(s9, Inches(4.2), Inches(5.5), Inches(2.2), "<₹400", "Hardware Cost")
    add_stat(s9, Inches(6.9), Inches(5.5), Inches(2.2), "8888", "UDP Port")
    add_stat(s9, Inches(9.6), Inches(5.5), Inches(2.2), "2s", "Broadcast Interval")

    # ══════════════ SLIDE 10: Use Cases ══════════════
    s10 = prs.slides.add_slide(blank_layout)
    add_header(s10, 10, "Impact & Roadmap", "Use Cases & Future Scope")

    uc = (
        "🎯 Real-World Use Cases\n\n"
        "• Disaster Relief — Communication when towers are down\n"
        "• Remote Expeditions — Trekking, camping, mountaineering\n"
        "• Event Management — Large venues with poor connectivity\n"
        "• Underground Facilities — Mines, bunkers, basements\n"
        "• Campus / Office — Private local-only communication"
    )
    add_text(s10, uc, Inches(1.5), Inches(2.5), Inches(5), Inches(4), 14, GRAY)

    fe = (
        "🚀 Future Enhancements\n\n"
        "• Multi-hop mesh relay (ESP32 to ESP32 chaining)\n"
        "• End-to-end message encryption (AES/RSA)\n"
        "• File / image transfer support\n"
        "• Group chat & broadcast channels\n"
        "• iOS client application"
    )
    add_text(s10, fe, Inches(7), Inches(2.5), Inches(5), Inches(4), 14, GRAY)

    # ══════════════ SLIDE 11: Thank You ══════════════
    s11 = prs.slides.add_slide(blank_layout)
    set_background(s11)
    add_text(
        s11,
        "11 / 11",
        Inches(11.5),
        Inches(0.3),
        Inches(1.5),
        Inches(0.5),
        12,
        CYAN,
        align=PP_ALIGN.RIGHT,
        font_name="Courier New",
    )
    add_text(s11, "📡", Inches(6.2), Inches(1.5), Inches(1), Inches(1), 40, align=PP_ALIGN.CENTER)
    add_text(
        s11,
        "END OF PRESENTATION",
        Inches(1),
        Inches(2.6),
        Inches(11.3),
        Inches(0.4),
        11,
        CYAN,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s11,
        "Thank You",
        Inches(1),
        Inches(3.0),
        Inches(11.3),
        Inches(1),
        54,
        WHITE,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        s11,
        "PeerLink — Connecting People, Offline.",
        Inches(1),
        Inches(4.0),
        Inches(11.3),
        Inches(0.5),
        16,
        GRAY,
        align=PP_ALIGN.CENTER,
    )

    add_stat(s11, Inches(3.5), Inches(5.0), Inches(2), "📱", "Android (Kotlin)")
    add_stat(s11, Inches(5.7), Inches(5.0), Inches(2), "⚙️", "ESP32 Firmware")
    add_stat(s11, Inches(7.9), Inches(5.0), Inches(2), "💬", "Offline Mesh")

    add_text(
        s11,
        "Built with Kotlin · Jetpack Compose · ESP32 Arduino · UDP Networking",
        Inches(1),
        Inches(6.8),
        Inches(11.3),
        Inches(0.5),
        11,
        RGBColor(100, 100, 100),
        align=PP_ALIGN.CENTER,
    )

    # Save Presentation
    prs.save("PeerLink_Presentation.pptx")
    print("Presentation generated successfully: PeerLink_Presentation.pptx")


if __name__ == "__main__":
    create_peerlink_presentation()

